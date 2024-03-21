import os
from pathlib import Path
from pptx import Presentation
import openpyxl
from docx import Document


def get_desktop_path():
    desktop = Path.home() / "Desktop"
    return str(desktop)


def create_pptx_file():
    desktop = get_desktop_path()
    prs = Presentation()


    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Новый слайд"


    file_path = os.path.join(desktop, "new_presentation.pptx")
    prs.save(file_path)
    print("Создан файл PowerPoint (pptx) на рабочем столе")


def create_xlsx_file():
    desktop = get_desktop_path()
    wb = openpyxl.Workbook()
    ws = wb.active


    ws['A1'] = 'Hello'
    ws['B1'] = 'World'


    file_path = os.path.join(desktop, "new_workbook.xlsx")
    wb.save(file_path)
    print("Создан файл Excel (xlsx) на рабочем столе")


def create_txt_file():
    desktop = get_desktop_path()
    file_path = os.path.join(desktop, "new_text_file.txt")
    with open(file_path, "w") as file:
        file.write("Hello, World!")
    print("Создан текстовый файл (txt) на рабочем столе")


def create_doc_file():
    desktop = get_desktop_path()
    doc = Document()
    doc.add_paragraph("Hello, World!")


    file_path = os.path.join(desktop, "new_document.docx")
    doc.save(file_path)
    print("Создан документ Word (docx) на рабочем столе")
